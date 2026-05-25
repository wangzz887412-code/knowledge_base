import os
from pathlib import Path
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation


def parse_file(file_path):
    file_path = Path(file_path)
    ext = file_path.suffix.lower()
    
    if ext == '.pdf':
        return parse_pdf(file_path)
    elif ext == '.docx':
        return parse_docx(file_path)
    elif ext == '.pptx':
        return parse_pptx(file_path)
    elif ext in ['.txt', '.md']:
        return parse_text(file_path)
    else:
        raise ValueError(f'Unsupported file type: {ext}')


def parse_pdf(file_path):
    text = ''
    reader = PdfReader(file_path)
    for page in reader.pages:
        text += page.extract_text() or ''
    return text


def parse_docx(file_path):
    doc = Document(file_path)
    text = ''
    for para in doc.paragraphs:
        text += para.text + '\n'
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text += cell.text + ' '
            text += '\n'
    return text


def parse_pptx(file_path):
    prs = Presentation(file_path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + '\n'
    return text


def parse_text(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()
